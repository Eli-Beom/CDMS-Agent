"""
CDMS-Agent 발표 슬라이드 5장
- 순백 배경 / 어두운 회색 글자
- 벡터 플랫 + 블루 악센트
- 픽토그램 · 인포그래픽
- 헤드라인 위치·크기 전 슬라이드 통일
- 마침표 없음
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

OUTPUT = r"C:\Users\SunbeomGwon\CDMS-Agent\pptx_work\CDMS_Agent_발표.pptx"

W = Inches(13.33)
H = Inches(7.50)

# ── 컬러 팔레트 ──────────────────────────────────────────────────────────────
BG       = RGBColor(0xFF, 0xFF, 0xFF)  # 순백
DARK     = RGBColor(0x1F, 0x29, 0x37)  # 헤드라인
BODY     = RGBColor(0x37, 0x41, 0x51)  # 본문
GRAY     = RGBColor(0x6B, 0x72, 0x80)  # 보조 텍스트
LGRAY    = RGBColor(0x9C, 0xA3, 0xAF)  # 연회색
BLUE     = RGBColor(0x25, 0x63, 0xEB)  # 블루 액센트
LBLUE    = RGBColor(0xEF, 0xF6, 0xFF)  # 연블루 배경
MBLUE    = RGBColor(0xDB, 0xEA, 0xFE)  # 중블루
BORDER   = RGBColor(0xE5, 0xE7, 0xEB)  # 카드 테두리
DBORDER  = RGBColor(0xD1, 0xD5, 0xDB)  # 짙은 테두리
RED      = RGBColor(0xDC, 0x26, 0x26)  # 경고
ORANGE   = RGBColor(0xD9, 0x77, 0x06)  # 주의
GREEN    = RGBColor(0x05, 0x96, 0x69)  # 성공
FONT     = "Pretendard"
FONT_ALT = "Apple SD Gothic Neo"
FONT_EN  = "Inter"

# ── 헤드라인 고정 좌표 (전 슬라이드 동일) ─────────────────────────────────
HL_X, HL_Y, HL_W, HL_H = Inches(0.55), Inches(0.32), Inches(12.23), Inches(0.62)
HL_SZ = 22       # pt
HL_COLOR = DARK

# ── helpers ──────────────────────────────────────────────────────────────────

def rgb_str(c: RGBColor) -> str:
    return str(c)

def _set_fill(spPr, color: RGBColor):
    for t in [qn("a:solidFill"), qn("a:noFill")]:
        for el in spPr.findall(t): spPr.remove(el)
    sf = etree.SubElement(spPr, qn("a:solidFill"))
    sc = etree.SubElement(sf, qn("a:srgbClr"))
    sc.set("val", rgb_str(color))
    _reorder_spPr(spPr)

def _set_no_fill(spPr):
    for t in [qn("a:solidFill"), qn("a:noFill")]:
        for el in spPr.findall(t): spPr.remove(el)
    etree.SubElement(spPr, qn("a:noFill"))

def _reorder_spPr(spPr):
    order = [qn("a:xfrm"), qn("a:prstGeom"), qn("a:noFill"),
             qn("a:solidFill"), qn("a:ln")]
    children = list(spPr)
    for c in children: spPr.remove(c)
    placed = []
    for tag in order:
        for c in children:
            if c.tag == tag and c not in placed:
                placed.append(c); break
    for c in children:
        if c not in placed: placed.append(c)
    for c in placed: spPr.append(c)

def _set_border(spPr, color: RGBColor = None, w=12700):
    ln = spPr.find(qn("a:ln"))
    if ln is None: ln = etree.SubElement(spPr, qn("a:ln"))
    for c in list(ln): ln.remove(c)
    if color:
        ln.set("w", str(w))
        sf = etree.SubElement(ln, qn("a:solidFill"))
        sc = etree.SubElement(sf, qn("a:srgbClr"))
        sc.set("val", rgb_str(color))
        pd = etree.SubElement(ln, qn("a:prstDash"))
        pd.set("val", "solid")
    else:
        etree.SubElement(ln, qn("a:noFill"))

def _set_radius(spPr, r):
    pg = spPr.find(qn("a:prstGeom"))
    if pg is None: return
    pg.set("prst", "roundRect")
    avl = pg.find(qn("a:avLst"))
    if avl is None: avl = etree.SubElement(pg, qn("a:avLst"))
    else:
        for c in list(avl): avl.remove(c)
    gd = etree.SubElement(avl, qn("a:gd"))
    gd.set("name", "adj"); gd.set("fmla", f"val {r}")

def rect(slide, x, y, w, h, fill=None, border=None, radius=0):
    sh = slide.shapes.add_shape(1, x, y, w, h)
    sp = sh.element; spPr = sp.find(qn("p:spPr"))
    if fill: _set_fill(spPr, fill)
    else: _set_no_fill(spPr)
    _set_border(spPr, border)
    if radius: _set_radius(spPr, radius)
    sh.text_frame.text = ""
    return sh

def arrow_right(slide, x, y, w, h, color=BLUE):
    sh = slide.shapes.add_shape(1, x, y, w, h)
    sp = sh.element; spPr = sp.find(qn("p:spPr"))
    pg = spPr.find(qn("a:prstGeom"))
    pg.set("prst", "rightArrow")
    avl = pg.find(qn("a:avLst"))
    if avl is None: avl = etree.SubElement(pg, qn("a:avLst"))
    else:
        for c in list(avl): avl.remove(c)
    for nm, vl in [("adj1","50000"),("adj2","50000")]:
        gd = etree.SubElement(avl, qn("a:gd"))
        gd.set("name", nm); gd.set("fmla", f"val {vl}")
    _set_fill(spPr, color); _set_border(spPr)
    sh.text_frame.text = ""; return sh

def arrow_down(slide, x, y, w, h, color=BLUE):
    sh = slide.shapes.add_shape(1, x, y, w, h)
    sp = sh.element; spPr = sp.find(qn("p:spPr"))
    pg = spPr.find(qn("a:prstGeom"))
    pg.set("prst", "downArrow")
    avl = pg.find(qn("a:avLst"))
    if avl is None: avl = etree.SubElement(pg, qn("a:avLst"))
    else:
        for c in list(avl): avl.remove(c)
    for nm, vl in [("adj1","50000"),("adj2","50000")]:
        gd = etree.SubElement(avl, qn("a:gd"))
        gd.set("name", nm); gd.set("fmla", f"val {vl}")
    _set_fill(spPr, color); _set_border(spPr)
    sh.text_frame.text = ""; return sh

def circle(slide, cx, cy, r, fill=BLUE, border=None):
    return rect(slide, cx - r, cy - r, r*2, r*2, fill=fill, border=border, radius=50000)

def txt(slide, x, y, w, h, paragraphs,
        anchor="middle", ml=Inches(0), mt=Inches(0),
        mr=Inches(0), mb=Inches(0), wrap=True):
    """
    paragraphs: list of (text, size_pt, bold, color, align)
    align: "l" | "c" | "r"
    """
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    body = tf._txBody
    bp = body.find(qn("a:bodyPr"))
    bp.set("lIns", str(int(ml))); bp.set("tIns", str(int(mt)))
    bp.set("rIns", str(int(mr))); bp.set("bIns", str(int(mb)))
    bp.set("anchor", {"top":"t","middle":"ctr","bottom":"b"}.get(anchor,"ctr"))

    for i, (text, sz, bold, color, align) in enumerate(paragraphs):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        pPr = para._p.get_or_add_pPr()
        para.alignment = {"l":PP_ALIGN.LEFT,"c":PP_ALIGN.CENTER,"r":PP_ALIGN.RIGHT}.get(align, PP_ALIGN.LEFT)
        etree.SubElement(pPr, qn("a:buNone"))
        if text == "":
            ep = etree.SubElement(para._p, qn("a:endParaRPr"))
            ep.set("lang","ko-KR"); ep.set("sz", str(int(sz*100))); continue
        run = para.add_run(); run.text = text
        rPr = run._r.get_or_add_rPr()
        rPr.set("lang","ko-KR"); rPr.set("sz", str(int(sz*100)))
        rPr.set("b","1" if bold else "0"); rPr.set("dirty","0")
        if color:
            sf = etree.SubElement(rPr, qn("a:solidFill"))
            sc = etree.SubElement(sf, qn("a:srgbClr"))
            sc.set("val", rgb_str(color))
        lat = etree.SubElement(rPr, qn("a:latin")); lat.set("typeface", FONT)
        ea  = etree.SubElement(rPr, qn("a:ea"));   ea.set("typeface", FONT)
    return tb

def headline(slide, text):
    """모든 슬라이드 헤드라인 — 위치·크기 고정"""
    # 블루 왼쪽 포인트 바
    rect(slide, HL_X - Inches(0.18), HL_Y + Inches(0.08),
         Inches(0.06), Inches(0.45), fill=BLUE)
    txt(slide, HL_X, HL_Y, HL_W, HL_H,
        [(text, HL_SZ, True, DARK, "l")],
        anchor="middle")
    # 구분선
    rect(slide, HL_X, HL_Y + HL_H + Inches(0.04),
         Inches(12.23), Inches(0.02), fill=BORDER)

def badge(slide, x, y, w, h, text, sz=9, fill=LBLUE, color=BLUE, radius=8000):
    rect(slide, x, y, w, h, fill=fill, radius=radius)
    txt(slide, x, y, w, h, [(text, sz, True, color, "c")], anchor="middle")

def icon_circle(slide, cx, cy, r, symbol, sym_sz=14, fill=BLUE, sym_color=None):
    circle(slide, cx, cy, r, fill=fill)
    txt(slide, cx-r, cy-r, r*2, r*2,
        [(symbol, sym_sz, True, sym_color or RGBColor(0xFF,0xFF,0xFF), "c")],
        anchor="middle")

def card(slide, x, y, w, h, fill=None, border=BORDER, radius=10000):
    return rect(slide, x, y, w, h,
                fill=fill or RGBColor(0xFF,0xFF,0xFF),
                border=border, radius=radius)

CONTENT_Y = Inches(1.18)  # 헤드라인 아래 콘텐츠 시작 Y


# ════════════════════════════════════════════════════════════════════════════════
# Slide 1 — 문제 정의
# ════════════════════════════════════════════════════════════════════════════════
def slide1(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])

    headline(sl, "DM 담당자는 매일 같은 작업을 손으로 반복한다")

    # ── 상단 스탯 3개 ──────────────────────────────────────────────────────
    stats = [
        ("5~10분",  "케이스 1건당 수동 입력 소요 시간",  RED),
        ("수십 건", "하루 처리해야 할 DVS 검증 건수",    ORANGE),
        ("100%",   "인적 오류가 개입될 수 있는 가능성",  BLUE),
    ]
    sw = Inches(3.6); sy = CONTENT_Y; sh = Inches(1.55)
    for i, (num, label, col) in enumerate(stats):
        sx = Inches(0.55) + i * Inches(4.29)
        card(sl, sx, sy, sw, sh, border=BORDER)
        # 컬러 상단 바
        rect(sl, sx, sy, sw, Inches(0.07), fill=col, radius=0)
        txt(sl, sx, sy + Inches(0.15), sw, Inches(0.75),
            [(num, 34, True, col, "c")], anchor="middle")
        txt(sl, sx, sy + Inches(0.92), sw, Inches(0.5),
            [(label, 9.5, False, GRAY, "c")], anchor="middle")

    # ── 문제 카드 3개 ──────────────────────────────────────────────────────
    problems = [
        ("🔁", "반복 클릭 작업",
         "CRF 페이지마다 필드 클릭 → 값 입력 → 저장 버튼 클릭을\n매 케이스, 매 방문마다 수행",
         RED, RGBColor(0xFF,0xF1,0xF1)),
        ("⚠", "인적 오류 리스크",
         "복사·붙여넣기, 오기록, 필드 혼동으로 인해\n데이터 품질이 담당자 역량에 의존",
         ORANGE, RGBColor(0xFF,0xF8,0xED)),
        ("⏱", "검증 시간 낭비",
         "DVS 시나리오 확인에 수시간 소요\n빠른 피드백 루프가 불가능한 구조",
         BLUE, LBLUE),
    ]
    py = CONTENT_Y + Inches(1.72); ph = Inches(3.45)
    pw = Inches(3.6)
    for i, (icon, title, desc, col, bg) in enumerate(problems):
        px = Inches(0.55) + i * Inches(4.29)
        card(sl, px, py, pw, ph, fill=bg, border=DBORDER)
        # 아이콘 원
        icon_circle(sl, px + Inches(0.5), py + Inches(0.55),
                    Inches(0.32), icon, sym_sz=16, fill=col)
        txt(sl, px + Inches(0.15), py + Inches(1.02), pw - Inches(0.3), Inches(0.4),
            [(title, 13, True, DARK, "l")], anchor="middle")
        txt(sl, px + Inches(0.15), py + Inches(1.48), pw - Inches(0.3), Inches(1.75),
            [(desc, 9.5, False, BODY, "l")], anchor="top",
            ml=Inches(0), mt=Inches(0))

    # ── 하단 결론 배너 ─────────────────────────────────────────────────────
    by = py + ph + Inches(0.2)
    rect(sl, Inches(0.55), by, Inches(12.23), Inches(0.48), fill=DARK, radius=8000)
    txt(sl, Inches(0.55), by, Inches(12.23), Inches(0.48),
        [("→  이 모든 작업을 Python 코드 한 줄로 대체할 수 없을까?", 11, True,
          RGBColor(0xFF,0xFF,0xFF), "c")], anchor="middle")


# ════════════════════════════════════════════════════════════════════════════════
# Slide 2 — 해결책: CDMS-Agent 개요
# ════════════════════════════════════════════════════════════════════════════════
def slide2(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])

    headline(sl, "CDMS-Agent: Python 코드로 CDMS 브라우저를 직접 조종")

    # ── 좌: 코드 블록 ──────────────────────────────────────────────────────
    cx = Inches(0.55); cw = Inches(5.8)
    cy = CONTENT_Y; ch = Inches(5.85)
    rect(sl, cx, cy, cw, ch,
         fill=RGBColor(0x1E,0x29,0x3B), radius=10000)

    # 코드 상단 dot bar
    for i, col in enumerate([RED, ORANGE, GREEN]):
        circle(sl, cx + Inches(0.28) + i*Inches(0.28),
               cy + Inches(0.28), Inches(0.07), fill=col)

    code_lines = [
        ("from cdm_agent_client import CDMSAgent", 9.5, False,
         RGBColor(0xA5,0xB4,0xFC)),
        ("", 8, False, None),
        ("agent = CDMSAgent(study_id='FAST_AF')", 9.5, False,
         RGBColor(0x86,0xEF,0xAC)),
        ("", 8, False, None),
        ("# 날짜 필드 자동 입력", 9, False, GRAY),
        ("agent.set_date('Visit date', '2025-01-15')", 9.5, False,
         RGBColor(0xFF,0xFF,0xFF)),
        ("", 8, False, None),
        ("# 저장 후 다음 페이지 이동", 9, False, GRAY),
        ("agent.click_save_next()", 9.5, False,
         RGBColor(0xFF,0xFF,0xFF)),
        ("", 8, False, None),
        ("# DVS 검증 결과 확인", 9, False, GRAY),
        ("result = agent.check_result('No Query')", 9.5, False,
         RGBColor(0xFF,0xFF,0xFF)),
        ("print(result)  # → 'PASS'", 9.5, False,
         RGBColor(0x86,0xEF,0xAC)),
    ]
    txt(sl, cx + Inches(0.2), cy + Inches(0.52),
        cw - Inches(0.4), ch - Inches(0.65),
        [(t, sz, b, c, "l") for t, sz, b, c in code_lines],
        anchor="top", ml=Inches(0), mt=Inches(0), wrap=False)

    # ── 우: 기능 카드 4개 ──────────────────────────────────────────────────
    features = [
        ("✏️", "CRF 필드 자동 입력",
         "날짜 · 텍스트 · 라디오 · 드롭다운\n모든 필드 타입 지원", BLUE),
        ("✅", "Query 자동 검증",
         "check_result()로 PASS / FAIL\n즉시 판정 후 결과 반환", GREEN),
        ("📄", "결과 리포트 자동화",
         "Excel · HTML 리포트 자동 생성\n원본 파일 덮어쓰기 없음", ORANGE),
        ("🔗", "Jupyter 완전 지원",
         "HTML 카드로 결과 즉시 시각화\n재현 가능한 스크립트 저장", RGBColor(0x71,0x48,0xF8)),
    ]
    fx = Inches(6.65); fw = Inches(6.12); fh = Inches(1.3)
    for i, (icon, title, desc, col) in enumerate(features):
        fy = CONTENT_Y + i * (fh + Inches(0.18))
        card(sl, fx, fy, fw, fh, border=BORDER)
        # 왼쪽 컬러 바
        rect(sl, fx, fy, Inches(0.07), fh, fill=col)
        icon_circle(sl, fx + Inches(0.52), fy + fh/2,
                    Inches(0.3), icon, sym_sz=14, fill=col)
        txt(sl, fx + Inches(0.95), fy + Inches(0.1),
            fw - Inches(1.05), Inches(0.45),
            [(title, 12, True, DARK, "l")], anchor="middle")
        txt(sl, fx + Inches(0.95), fy + Inches(0.55),
            fw - Inches(1.05), Inches(0.65),
            [(desc, 9, False, BODY, "l")], anchor="top",
            ml=Inches(0), mt=Inches(0))


# ════════════════════════════════════════════════════════════════════════════════
# Slide 3 — 아키텍처: WebSocket 시각화
# ════════════════════════════════════════════════════════════════════════════════
def slide3(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])

    headline(sl, "4계층 구조: Python 명령이 CDMS DOM까지 전달되는 흐름")

    # ── 수직 플로우 다이어그램 ─────────────────────────────────────────────
    layers = [
        ("① Python / Jupyter Notebook",
         "pip install cdm-agent-client\nagent.set_date(...)",
         BLUE, LBLUE, "HTTP POST  /api/cdm-agent/run-case"),
        ("② CDM Agent Daemon  (Node.js · Port 3100)",
         "로컬 실행 서버\nHTTP 요청을 수신해 WebSocket으로 변환",
         DARK, RGBColor(0xF1,0xF5,0xF9), "WebSocket  ws://127.0.0.1:3100/ws/cdm-agent"),
        ("③ Chrome Extension  (MV3 Service Worker)",
         "WebSocket 클라이언트\nchromeScripting.executeScript (world: MAIN)",
         RGBColor(0x71,0x48,0xF8), RGBColor(0xF5,0xF3,0xFF), "executeScript → MAIN world"),
        ("④ CDMS Page DOM  ←  browser-runner-core.js",
         "React DOM 직접 조작\nfindRow → setNativeValue → 이벤트 발생",
         GREEN, RGBColor(0xEC,0xFD,0xF5), None),
    ]

    bw = Inches(5.6); bh = Inches(1.1)
    bx = Inches(0.75)
    by_start = CONTENT_Y + Inches(0.1)
    gap_box   = Inches(0.38)

    for i, (title, desc, col, bg, proto) in enumerate(layers):
        by = by_start + i * (bh + gap_box)

        # 메인 박스
        card(sl, bx, by, bw, bh, fill=bg, border=DBORDER, radius=8000)
        rect(sl, bx, by, Inches(0.09), bh, fill=col, radius=0)

        txt(sl, bx + Inches(0.22), by + Inches(0.1),
            bw - Inches(0.3), Inches(0.42),
            [(title, 11, True, col if bg != LBLUE else DARK, "l")],
            anchor="middle")
        txt(sl, bx + Inches(0.22), by + Inches(0.52),
            bw - Inches(0.3), Inches(0.5),
            [(desc, 8.5, False, BODY, "l")], anchor="top",
            ml=Inches(0), mt=Inches(0))

        # 프로토콜 화살표 + 라벨
        if i < len(layers) - 1:
            ay = by + bh
            arrow_down(sl, bx + bw/2 - Inches(0.12),
                       ay + Inches(0.05), Inches(0.24), Inches(0.25))
            if proto:
                # WebSocket 구간은 강조
                is_ws = "WebSocket" in proto
                proto_fill = BLUE if is_ws else BORDER
                proto_col  = RGBColor(0xFF,0xFF,0xFF) if is_ws else GRAY
                pw = Inches(3.2)
                px = bx + bw/2 - pw/2
                badge(sl, px + Inches(1.45), ay + Inches(0.07),
                      pw, Inches(0.22),
                      proto, sz=8,
                      fill=proto_fill if is_ws else LBLUE,
                      color=proto_col if is_ws else BLUE,
                      radius=5000)

    # ── 우측: 상세 설명 패널 ──────────────────────────────────────────────
    px = Inches(7.1); pw = Inches(5.75)

    # WebSocket 핵심 강조 박스
    wy = CONTENT_Y + Inches(0.1)
    rect(sl, px, wy, pw, Inches(2.05), fill=LBLUE, border=BLUE, radius=10000)
    txt(sl, px + Inches(0.2), wy + Inches(0.15),
        pw - Inches(0.4), Inches(0.42),
        [("WebSocket — 핵심 연결 채널", 12, True, BLUE, "l")], anchor="middle")
    ws_details = [
        "• Daemon(서버)에 Extension(클라이언트)이 접속",
        "• 상시 연결 유지 — 끊기면 자동 재연결",
        "• 명령 수신 → executeScript 호출 → 결과 반송",
        "• 라디오 버튼은 CDP(isTrusted) 별도 처리",
    ]
    for j, line in enumerate(ws_details):
        txt(sl, px + Inches(0.2),
            wy + Inches(0.62) + j * Inches(0.34),
            pw - Inches(0.4), Inches(0.34),
            [(line, 9, False, BODY, "l")], anchor="middle",
            ml=Inches(0), mt=Inches(0))

    # 실행 흐름 요약 박스
    ry = wy + Inches(2.22)
    rect(sl, px, ry, pw, Inches(3.0), fill=RGBColor(0xFF,0xFF,0xFF),
         border=BORDER, radius=10000)
    txt(sl, px + Inches(0.2), ry + Inches(0.15),
        pw - Inches(0.4), Inches(0.38),
        [("실행 흐름 요약", 11, True, DARK, "l")], anchor="middle")
    steps = [
        ("1", "agent.set_date() 호출",         BLUE),
        ("2", "HTTP POST → Daemon 수신",         DARK),
        ("3", "WebSocket 'command' 전송",         BLUE),
        ("4", "executeScript → runner.run()",  RGBColor(0x71,0x48,0xF8)),
        ("5", "DOM 조작 → StepResult 반환",     GREEN),
    ]
    for j, (num, step, col) in enumerate(steps):
        sy = ry + Inches(0.6) + j * Inches(0.46)
        circle(sl, px + Inches(0.35), sy + Inches(0.14),
               Inches(0.14), fill=col)
        txt(sl, px + Inches(0.1), sy + Inches(0.02),
            Inches(0.5), Inches(0.28),
            [(num, 8, True, RGBColor(0xFF,0xFF,0xFF), "c")], anchor="middle")
        txt(sl, px + Inches(0.62), sy,
            pw - Inches(0.75), Inches(0.3),
            [(step, 9.5, False, BODY, "l")], anchor="middle",
            ml=Inches(0), mt=Inches(0))


# ════════════════════════════════════════════════════════════════════════════════
# Slide 4 — Python 패키지 계층도
# ════════════════════════════════════════════════════════════════════════════════
def slide4(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])

    headline(sl, "cdm_agent_client: 3개 모듈, 하나의 통합 인터페이스")

    # ── 좌: 트리 다이어그램 ────────────────────────────────────────────────
    tx = Inches(0.55); tw = Inches(5.9)
    ty = CONTENT_Y + Inches(0.1)

    # 루트 패키지
    rect(sl, tx, ty, tw, Inches(0.58), fill=BLUE, radius=8000)
    txt(sl, tx, ty, tw, Inches(0.58),
        [("📦  cdm_agent_client/", 12, True,
          RGBColor(0xFF,0xFF,0xFF), "l")],
        anchor="middle", ml=Inches(0.2))

    modules = [
        ("__init__.py",   "공개 인터페이스 통합 진입점",  DARK,  LBLUE,
         ["CDMSAgent, CDMAgent", "PageSnapshot, StepResult",
          "CDMAgentError 계열"]),
        ("client.py",     "CDMSAgent 메인 클래스",         BLUE,  LBLUE,
         ["inspect / list_pages / ping",
          "set_date / set_text / select_*",
          "click_save / go_to_page / go_back",
          "has_query / check_result"]),
        ("models.py",     "반환값 데이터 구조",            RGBColor(0x71,0x48,0xF8), LBLUE,
         ["PageSnapshot — 페이지 스냅샷",
          "PageList — CRF 페이지 목록",
          "StepResult — 조작 결과"]),
        ("exceptions.py", "에러 클래스 계층",              RED,   RGBColor(0xFF,0xF1,0xF1),
         ["CDMAgentError (부모)",
          "DaemonNotRunningError",
          "NoBrowserClientError",
          "StepFailedError"]),
    ]

    mh_list = [0.58, 1.45, 1.12, 1.12]
    cur_y = ty + Inches(0.75)
    INDENT = Inches(0.35)
    for idx, (fname, role, col, bg, items) in enumerate(modules):
        mh = Inches(mh_list[idx])
        # 연결선
        rect(sl, tx + INDENT - Inches(0.02), cur_y,
             Inches(0.02), mh, fill=BORDER)
        rect(sl, tx + INDENT, cur_y + mh/2 - Inches(0.01),
             Inches(0.2), Inches(0.02), fill=BORDER)

        mx = tx + INDENT + Inches(0.2)
        mw = tw - INDENT - Inches(0.2)
        card(sl, mx, cur_y, mw, mh, fill=bg, border=DBORDER, radius=6000)
        rect(sl, mx, cur_y, Inches(0.06), mh, fill=col, radius=0)

        txt(sl, mx + Inches(0.15), cur_y + Inches(0.05),
            mw - Inches(0.2), Inches(0.38),
            [(fname, 10, True, col, "l")], anchor="middle")
        if mh > Inches(0.65):
            txt(sl, mx + Inches(0.15), cur_y + Inches(0.4),
                mw - Inches(0.2), mh - Inches(0.45),
                [(it, 8, False, BODY, "l") for it in items],
                anchor="top", ml=Inches(0), mt=Inches(0))
        cur_y += mh + Inches(0.12)

    # ── 우: CDMSAgent 메서드 그룹 상세 ────────────────────────────────────
    rx = Inches(6.75); rw = Inches(6.12)

    groups = [
        ("📋  페이지 정보",
         [("inspect(client_id?)",  "→ PageSnapshot  현재 페이지 스냅샷"),
          ("list_pages()",         "→ PageList  사이드바 CRF 목록"),
          ("ping()",               "→ bool  Daemon 연결 확인"),
          ("clients()",            "→ list  연결된 Extension 목록")],
         BLUE),
        ("✏️  필드 입력",
         [("set_date(row_label, value)",    "날짜 필드 (캘린더 팝업 포함)"),
          ("set_text(row_label, value)",    "텍스트 · 숫자 필드"),
          ("select_option(row, option)",    "드롭다운 선택"),
          ("select_radio(row, option)",     "라디오 버튼 (CDP 처리)")],
         RGBColor(0x71,0x48,0xF8)),
        ("🖱️  버튼 · 네비게이션",
         [("click_save / click_save_next()", "Save / Save & Next 클릭"),
          ("go_to_page(segment)",           "\"V2/DM\" → URL 세그먼트 이동"),
          ("go_back() / navigate_to(url)",  "브라우저 히스토리 / 직접 이동")],
         GREEN),
        ("✅  검증",
         [("has_query()",          "→ bool  Query 존재 여부"),
          ("check_result(expected)","→ PASS | FAIL  DVS 검증")],
         ORANGE),
    ]

    gh_list = [1.32, 1.47, 1.18, 0.85]
    cur_gy = CONTENT_Y + Inches(0.1)
    for gi, (gtitle, methods, col) in enumerate(groups):
        gh = Inches(gh_list[gi])
        card(sl, rx, cur_gy, rw, gh, border=BORDER, radius=8000)
        rect(sl, rx, cur_gy, Inches(0.07), gh, fill=col, radius=0)
        txt(sl, rx + Inches(0.18), cur_gy + Inches(0.05),
            rw - Inches(0.25), Inches(0.36),
            [(gtitle, 10, True, col, "l")], anchor="middle")
        for mi, (mname, mdesc) in enumerate(methods):
            my = cur_gy + Inches(0.42) + mi * Inches(0.2)
            txt(sl, rx + Inches(0.18), my, Inches(2.85), Inches(0.2),
                [(mname, 8, True, DARK, "l")], anchor="middle",
                ml=Inches(0), mt=Inches(0))
            txt(sl, rx + Inches(3.1), my, rw - Inches(3.2), Inches(0.2),
                [(mdesc, 8, False, GRAY, "l")], anchor="middle",
                ml=Inches(0), mt=Inches(0))
        cur_gy += gh + Inches(0.14)


# ════════════════════════════════════════════════════════════════════════════════
# Slide 5 — 기대 효과
# ════════════════════════════════════════════════════════════════════════════════
def slide5(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])

    headline(sl, "자동화가 가져오는 DM 업무의 3가지 변화")

    changes = [
        ("⏱", "처리 속도",
         "케이스당 10분", "수십 초",
         "반복 클릭·저장 과정을 스크립트가\n대신 수행해 처리 속도가 비약적으로 향상",
         RED, GREEN),
        ("🎯", "데이터 품질",
         "인적 오류 발생", "자동 검증 통과",
         "check_result()가 입력 즉시 Query 발생\n여부를 확인해 오류를 사전에 차단",
         ORANGE, GREEN),
        ("🔄", "재현성",
         "매번 수동 반복", "코드로 언제든 재현",
         "Jupyter 스크립트로 저장해두면\n동일 작업을 언제든 동일하게 재실행 가능",
         GRAY, BLUE),
    ]

    cw = Inches(3.95); ch = Inches(4.65)
    cy = CONTENT_Y + Inches(0.15)
    for i, (icon, cat, before, after, desc, bcol, acol) in enumerate(changes):
        cx = Inches(0.55) + i * Inches(4.32)

        # 메인 카드
        card(sl, cx, cy, cw, ch, border=BORDER, radius=12000)

        # 아이콘 + 카테고리
        icon_circle(sl, cx + cw/2, cy + Inches(0.58),
                    Inches(0.38), icon, sym_sz=20, fill=BLUE)
        txt(sl, cx, cy + Inches(1.08), cw, Inches(0.38),
            [(cat, 13, True, DARK, "c")], anchor="middle")

        # Before / After
        bw = cw / 2 - Inches(0.12)
        # Before
        rect(sl, cx + Inches(0.1), cy + Inches(1.58),
             bw, Inches(0.72), fill=RGBColor(0xFF,0xF1,0xF1), radius=6000)
        txt(sl, cx + Inches(0.1), cy + Inches(1.58),
            bw, Inches(0.24),
            [("BEFORE", 7, True, bcol, "c")], anchor="middle")
        txt(sl, cx + Inches(0.1), cy + Inches(1.82),
            bw, Inches(0.46),
            [(before, 9, True, bcol, "c")], anchor="middle")
        # Arrow
        arrow_right(sl, cx + cw/2 - Inches(0.14),
                    cy + Inches(1.78), Inches(0.28), Inches(0.24), color=BLUE)
        # After
        rect(sl, cx + cw/2 + Inches(0.04), cy + Inches(1.58),
             bw, Inches(0.72), fill=RGBColor(0xEC,0xFD,0xF5), radius=6000)
        txt(sl, cx + cw/2 + Inches(0.04), cy + Inches(1.58),
            bw, Inches(0.24),
            [("AFTER", 7, True, acol, "c")], anchor="middle")
        txt(sl, cx + cw/2 + Inches(0.04), cy + Inches(1.82),
            bw, Inches(0.46),
            [(after, 9, True, acol, "c")], anchor="middle")

        # 구분선
        rect(sl, cx + Inches(0.15), cy + Inches(2.42),
             cw - Inches(0.3), Inches(0.02), fill=BORDER)

        # 설명
        txt(sl, cx + Inches(0.15), cy + Inches(2.52),
            cw - Inches(0.3), Inches(1.85),
            [(desc, 9, False, BODY, "c")], anchor="top",
            ml=Inches(0), mt=Inches(0))

    # ── 하단 로드맵 배너 ───────────────────────────────────────────────────
    ry = cy + ch + Inches(0.22)
    rect(sl, Inches(0.55), ry, Inches(12.23), Inches(0.55),
         fill=LBLUE, border=BLUE, radius=8000)
    txt(sl, Inches(0.55), ry, Inches(12.23), Inches(0.55),
        [("향후: 멀티 스터디 배치 실행  ·  AI 기반 CRF 이상 탐지  ·  결과 대시보드 자동화", 10, False, BLUE, "c")],
        anchor="middle")


# ════════════════════════════════════════════════════════════════════════════════
# 빌드
# ════════════════════════════════════════════════════════════════════════════════
def main():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    slide1(prs)
    slide2(prs)
    slide3(prs)
    slide4(prs)
    slide5(prs)

    if os.path.exists(OUTPUT): os.remove(OUTPUT)
    prs.save(OUTPUT)
    size = os.path.getsize(OUTPUT) / 1024
    print(f"저장 완료: {OUTPUT}  ({size:.0f} KB)  |  {len(prs.slides)}장")

if __name__ == "__main__":
    main()
