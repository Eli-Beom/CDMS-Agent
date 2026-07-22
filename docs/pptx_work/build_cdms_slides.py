"""
CDMS-Agent 소개 슬라이드 3장 생성
- 템플릿 불필요, 완전히 새 PPTX로 생성
- JNPMEDI 색상/폰트 직접 적용
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import copy, os

OUTPUT = r"C:\Users\SunbeomGwon\CDMS-Agent\pptx_work\CDMS_Agent_소개_슬라이드.pptx"

# ── 색상 ──────────────────────────────────────────────────────────────────────
C_NAVY   = RGBColor(0x0B, 0x1F, 0x3B)
C_BLUE   = RGBColor(0x1E, 0x4F, 0xA8)
C_LBLUE  = RGBColor(0xEF, 0xF6, 0xFF)
C_CARD   = RGBColor(0xF0, 0xF4, 0xF8)
C_BORDER = RGBColor(0xE5, 0xE7, 0xEB)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK   = RGBColor(0x1A, 0x2B, 0x4A)
C_GRAY   = RGBColor(0x6B, 0x72, 0x80)
C_ARROW  = RGBColor(0x3B, 0x82, 0xF6)
C_GREEN  = RGBColor(0x05, 0x96, 0x69)
C_ORANGE = RGBColor(0xD9, 0x77, 0x06)
C_YELLOW = RGBColor(0xFF, 0xFF, 0x00)
C_TEAL   = RGBColor(0x15, 0x65, 0xC0)

W = Inches(13.33)
H = Inches(7.50)
FONT = "Noto Sans KR"

# ── 헬퍼 ──────────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill=None, border=None, radius=0):
    """사각형 or 둥근 사각형 추가 (텍스트 없음)."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE (auto_shape_type 1 = rounded_rect handled via XML)
        x, y, w, h
    )
    sp = shape.element
    # 채우기
    spPr = sp.find(qn("p:spPr"))
    # prstGeom 변경
    prstGeom = spPr.find(qn("a:prstGeom"))
    if radius and prstGeom is not None:
        prstGeom.set("prst", "roundRect")
        avLst = prstGeom.find(qn("a:avLst"))
        if avLst is None:
            avLst = etree.SubElement(prstGeom, qn("a:avLst"))
        else:
            for c in list(avLst): avLst.remove(c)
        gd = etree.SubElement(avLst, qn("a:gd"))
        gd.set("name", "adj"); gd.set("fmla", f"val {radius}")
    # solidFill
    if fill is not None:
        _set_fill(spPr, fill)
    else:
        _set_no_fill(spPr)
    # 테두리
    ln = spPr.find(qn("a:ln"))
    if ln is None:
        ln = etree.SubElement(spPr, qn("a:ln"))
    if border is not None:
        ln.set("w", "12700")
        sf = etree.SubElement(ln, qn("a:solidFill"))
        sc = etree.SubElement(sf, qn("a:srgbClr"))
        sc.set("val", str(border))
        pd = etree.SubElement(ln, qn("a:prstDash"))
        pd.set("val", "solid")
    else:
        for c in list(ln): ln.remove(c)
        etree.SubElement(ln, qn("a:noFill"))
    shape.text_frame.text = ""
    return shape


def _set_fill(spPr, color):
    for old in spPr.findall(qn("a:solidFill")) + spPr.findall(qn("a:noFill")):
        spPr.remove(old)
    sf = etree.SubElement(spPr, qn("a:solidFill"))
    sc = etree.SubElement(sf, qn("a:srgbClr"))
    sc.set("val", str(color))
    # move fill before prstGeom? no – just insert at correct position
    # reorder: noFill/solidFill should come after xfrm and prstGeom
    children = list(spPr)
    for c in children:
        spPr.remove(c)
    order = [qn("a:xfrm"), qn("a:prstGeom"), qn("a:solidFill"), qn("a:ln")]
    placed = []
    for tag in order:
        for c in children:
            if c.tag == tag:
                placed.append(c); break
    for c in children:
        if c not in placed:
            placed.append(c)
    for c in placed:
        spPr.append(c)


def _set_no_fill(spPr):
    for old in spPr.findall(qn("a:solidFill")) + spPr.findall(qn("a:noFill")):
        spPr.remove(old)
    etree.SubElement(spPr, qn("a:noFill"))


def add_textbox(slide, x, y, w, h,
                paragraphs,          # list of (text, sz_pt, bold, color, align, lang)
                vert_anchor="top",
                word_wrap=True,
                margin_l=Inches(0.05), margin_t=Inches(0.03),
                margin_r=Inches(0.05), margin_b=Inches(0.03)):
    """텍스트박스 추가."""
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf = txb.text_frame
    tf.word_wrap = word_wrap

    # body properties
    txBody = tf._txBody
    bodyPr = txBody.find(qn("a:bodyPr"))
    bodyPr.set("lIns", str(int(margin_l)))
    bodyPr.set("tIns", str(int(margin_t)))
    bodyPr.set("rIns", str(int(margin_r)))
    bodyPr.set("bIns", str(int(margin_b)))
    anchor_map = {"top": "t", "middle": "ctr", "bottom": "b"}
    bodyPr.set("anchor", anchor_map.get(vert_anchor, "t"))

    for i, (text, sz, bold, color, align, lang) in enumerate(paragraphs):
        if i == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()

        pPr = para._p.get_or_add_pPr()
        # alignment
        align_map = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER, "r": PP_ALIGN.RIGHT}
        para.alignment = align_map.get(align, PP_ALIGN.LEFT)
        # no bullet
        buNone = etree.SubElement(pPr, qn("a:buNone"))

        if text == "":
            # empty spacer paragraph
            endPr = etree.SubElement(para._p, qn("a:endParaRPr"))
            endPr.set("lang", lang)
            endPr.set("sz", str(int(sz * 100)))
            continue

        run = para.add_run()
        run.text = text
        rPr = run._r.get_or_add_rPr()
        rPr.set("lang", lang)
        rPr.set("sz", str(int(sz * 100)))
        rPr.set("b", "1" if bold else "0")
        rPr.set("kern", "0")
        rPr.set("dirty", "0")
        # color
        if color is not None:
            sf = etree.SubElement(rPr, qn("a:solidFill"))
            sc = etree.SubElement(sf, qn("a:srgbClr"))
            sc.set("val", str(color))
        # font
        lat = etree.SubElement(rPr, qn("a:latin"))
        lat.set("typeface", FONT)
        ea = etree.SubElement(rPr, qn("a:ea"))
        ea.set("typeface", FONT)

    return txb


def add_arrow_down(slide, x, y, w, h, color=C_ARROW):
    """아래 방향 화살표."""
    shape = slide.shapes.add_shape(1, x, y, w, h)
    sp = shape.element
    spPr = sp.find(qn("p:spPr"))
    prstGeom = spPr.find(qn("a:prstGeom"))
    prstGeom.set("prst", "downArrow")
    avLst = prstGeom.find(qn("a:avLst"))
    if avLst is None:
        avLst = etree.SubElement(prstGeom, qn("a:avLst"))
    else:
        for c in list(avLst): avLst.remove(c)
    for name, val in [("adj1", "50000"), ("adj2", "50000")]:
        gd = etree.SubElement(avLst, qn("a:gd"))
        gd.set("name", name); gd.set("fmla", f"val {val}")
    _set_fill(spPr, color)
    ln = spPr.find(qn("a:ln"))
    if ln is None: ln = etree.SubElement(spPr, qn("a:ln"))
    for c in list(ln): ln.remove(c)
    etree.SubElement(ln, qn("a:noFill"))
    shape.text_frame.text = ""
    return shape


# ── 복합 컴포넌트 ──────────────────────────────────────────────────────────────

def draw_title_bar(slide, title, subtitle=None):
    """상단 네이비 타이틀 바."""
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.95), fill=C_NAVY)
    paras = [(title, 18, True, C_WHITE, "l", "ko-KR")]
    if subtitle:
        paras.append((subtitle, 11, False, RGBColor(0xCA, 0xDC, 0xFC), "l", "ko-KR"))
    add_textbox(slide, Inches(0.47), Inches(0.07), Inches(12.0), Inches(0.82),
                paras, vert_anchor="middle",
                margin_l=Inches(0), margin_t=Inches(0),
                margin_r=Inches(0), margin_b=Inches(0))
    # Confidential
    add_textbox(slide, Inches(11.45), Inches(0.21), Inches(1.38), Inches(0.35),
                [("Confidential", 11, True, C_YELLOW, "r", "en-US")],
                vert_anchor="middle",
                margin_l=Inches(0), margin_t=Inches(0),
                margin_r=Inches(0), margin_b=Inches(0))


def draw_card(slide, x, y, w, h, num, label, body_paragraphs):
    """
    카드: 둥근 사각형 배경 + 네이비 왼쪽 바 + 번호 원 + 레이블 + 본문
    body_paragraphs: list of (text, sz, bold, color, lang)
    """
    # 카드 배경
    add_rect(slide, x, y, w, h, fill=C_CARD, border=C_BORDER, radius=14815)
    # 왼쪽 accent bar
    add_rect(slide, x, y, Inches(0.07), h, fill=C_NAVY, radius=150000)
    # 번호 원
    cx, cy = x + Inches(0.12), y + Inches(0.10)
    add_rect(slide, cx, cy, Inches(0.28), Inches(0.28), fill=C_BLUE, radius=50000)
    add_textbox(slide, cx, cy, Inches(0.28), Inches(0.28),
                [(str(num), 9, True, C_WHITE, "c", "en-US")],
                vert_anchor="middle",
                margin_l=Inches(0), margin_t=Inches(0),
                margin_r=Inches(0), margin_b=Inches(0))
    # 레이블
    add_textbox(slide, x + Inches(0.45), y + Inches(0.08), w - Inches(0.52), Inches(0.32),
                [(label, 11, True, C_DARK, "l", "ko-KR")],
                vert_anchor="middle",
                margin_l=Inches(0), margin_t=Inches(0),
                margin_r=Inches(0), margin_b=Inches(0))
    # 본문
    paras = [(t, sz, b, col, "l", lang) for (t, sz, b, col, lang) in body_paragraphs]
    add_textbox(slide, x + Inches(0.45), y + Inches(0.43), w - Inches(0.55), h - Inches(0.53),
                paras, vert_anchor="top",
                margin_l=Inches(0), margin_t=Inches(0),
                margin_r=Inches(0.05), margin_b=Inches(0))


def draw_arch_box(slide, x, y, w, h, text, fill=C_BLUE, sz=9):
    add_rect(slide, x, y, w, h, fill=fill, radius=8000)
    add_textbox(slide, x + Inches(0.05), y, w - Inches(0.1), h,
                [(text, sz, True, C_WHITE, "c", "ko-KR")],
                vert_anchor="middle",
                margin_l=Inches(0), margin_t=Inches(0),
                margin_r=Inches(0), margin_b=Inches(0))


def draw_tag(slide, x, y, w, h, text):
    add_rect(slide, x, y, w, h, fill=C_LBLUE, radius=12000)
    add_textbox(slide, x, y, w, h,
                [(text, 8, False, C_BLUE, "c", "en-US")],
                vert_anchor="middle",
                margin_l=Inches(0.04), margin_t=Inches(0),
                margin_r=Inches(0.04), margin_b=Inches(0))


def draw_section_header(slide, x, y, w, text):
    add_textbox(slide, x, y, w, Inches(0.36),
                [(text, 11, True, C_NAVY, "l", "ko-KR")],
                vert_anchor="middle",
                margin_l=Inches(0), margin_t=Inches(0),
                margin_r=Inches(0), margin_b=Inches(0))


# ── 슬라이드 1: 아키텍처 ──────────────────────────────────────────────────────

def build_slide1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    slide.shapes.title  # may not exist on blank

    draw_title_bar(slide,
                   "CDMS-Agent 아키텍처 및 작동 원리",
                   "브라우저 자동화 4-Layer 구조")

    # ── 좌: 아키텍처 흐름도 ──
    draw_section_header(slide, Inches(0.45), Inches(1.05), Inches(5.9),
                        "작동 원리 (4-Layer Architecture)")

    bx = Inches(0.55)
    bw = Inches(5.35)
    bh = Inches(0.55)
    ay = [Inches(1.50), Inches(2.32), Inches(3.14), Inches(3.96)]
    arr_x = bx + Inches(2.5)
    arr_w = Inches(0.22)
    arr_h = Inches(0.28)

    boxes = [
        ("① Jupyter Notebook / Python Script", C_BLUE),
        ("② CDM Agent Daemon  (Node.js · Port 3100)", C_NAVY),
        ("③ Chrome Extension  (MV3 Service Worker)", C_TEAL),
        ("④ CDMS Page DOM  ←  browser-runner-core.js", RGBColor(0x05, 0x96, 0x69)),
    ]
    tags = ["pip install cdm-agent-client", "localhost:3100",
            "chrome.scripting API", "executeScript (MAIN world)"]
    protocols = ["HTTP (requests)", "WebSocket", "executeScript (MAIN world)"]

    for i, (text, fill) in enumerate(boxes):
        draw_arch_box(slide, bx, ay[i], bw, bh, text, fill=fill, sz=9)
        draw_tag(slide, bx + bw + Inches(0.08), ay[i] + Inches(0.13),
                 Inches(1.55), Inches(0.29), tags[i])
        if i < 3:
            add_arrow_down(slide, arr_x, ay[i] + bh, arr_w, arr_h)
            add_textbox(slide,
                        arr_x + arr_w + Inches(0.05),
                        ay[i] + bh + Inches(0.04),
                        Inches(1.5), Inches(0.20),
                        [(protocols[i], 8, False, C_GRAY, "l", "en-US")],
                        vert_anchor="middle",
                        margin_l=Inches(0), margin_t=Inches(0),
                        margin_r=Inches(0), margin_b=Inches(0))

    # ── 우: 설명 카드 3개 ──
    rx = Inches(7.05)
    rw = Inches(5.82)

    draw_section_header(slide, rx, Inches(1.05), rw, "레이어별 역할")

    draw_card(slide, rx, Inches(1.50), rw, Inches(1.72), 1,
              "Python 클라이언트 패키지",
              [
                  ("cdm_agent_client가 Python ↔ Daemon 레이어를 담당합니다.", 9, False, C_DARK, "ko-KR"),
                  ("", 8, False, C_DARK, "ko-KR"),
                  ("set_date / set_text / select_radio / select_option", 8.5, False, C_GRAY, "en-US"),
                  ("click_save_next / go_to_page / go_back / navigate_to", 8.5, False, C_GRAY, "en-US"),
                  ("inspect() → PageSnapshot  |  list_pages() → PageList", 8.5, False, C_GRAY, "en-US"),
                  ("결과를 Jupyter HTML 카드로 즉시 시각화", 8.5, False, C_GRAY, "ko-KR"),
              ])

    draw_card(slide, rx, Inches(3.30), rw, Inches(1.72), 2,
              "CDM Agent Daemon & Chrome Extension",
              [
                  ("로컬 Node.js 서버가 Python 요청을 받아 WebSocket으로 전달합니다.", 9, False, C_DARK, "ko-KR"),
                  ("", 8, False, C_DARK, "ko-KR"),
                  ("Daemon: cdm-agent-platform (별도 저장소)", 8.5, False, C_GRAY, "ko-KR"),
                  ("Extension: MV3 Service Worker 방식으로 Chrome에 로드", 8.5, False, C_GRAY, "ko-KR"),
                  ("MAIN world 스크립트 실행 → CRF 페이지 DOM 직접 접근", 8.5, False, C_GRAY, "ko-KR"),
              ])

    draw_card(slide, rx, Inches(5.10), rw, Inches(1.72), 3,
              "CDMS Page DOM 조작 결과",
              [
                  ("browser-runner-core.js가 CRF DOM을 직접 조작해 자동 처리합니다.", 9, False, C_DARK, "ko-KR"),
                  ("", 8, False, C_DARK, "ko-KR"),
                  ("날짜 · 텍스트 · 라디오 · 드롭다운 입력 / Save & Next 클릭", 8.5, False, C_GRAY, "ko-KR"),
                  ("check_result() → Query 발생 여부 즉시 검증", 8.5, False, C_GRAY, "ko-KR"),
                  ("실행 결과: passed / failed / blocked / skipped 반환", 8.5, False, C_GRAY, "en-US"),
              ])


# ── 슬라이드 2: 패키지 구조 ───────────────────────────────────────────────────

def build_slide2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    draw_title_bar(slide,
                   "cdm_agent_client 패키지 구조 및 작동 방식",
                   "CDMAgent (브라우저 조작) + DVSRunner (DVS 검증 오케스트레이터)")

    # ── 좌: 패키지 트리 ──
    draw_section_header(slide, Inches(0.42), Inches(1.05), Inches(6.0),
                        "패키지 구조  (src/cdm_agent_client/)")

    tree_lines = [
        ("src/cdm_agent_client/",                                   9.5, True,  C_NAVY, "en-US"),
        ("├── __init__.py",                                         9,   False, C_GRAY, "en-US"),
        ("├── client.py       ← CDMAgent 클래스",                   9,   False, C_BLUE, "ko-KR"),
        ("├── models.py       ← PageSnapshot · PageList · StepResult", 9, False, C_BLUE, "en-US"),
        ("├── exceptions.py   ← DaemonNotRunning · NoBrowser · StepFailed", 9, False, C_BLUE, "en-US"),
        ("└── dvs/",                                                9,   False, C_NAVY, "en-US"),
        ("    ├── runner.py   ← DVSRunner (핵심 오케스트레이터)",    9,   False, C_BLUE, "ko-KR"),
        ("    ├── schema.py   ← DVSRow · DVSResult · Precondition", 9,   False, C_BLUE, "en-US"),
        ("    ├── parser.py   ← Excel → DVSRow 파싱",               9,   False, C_BLUE, "ko-KR"),
        ("    ├── planner.py  ← Precondition → ActionStep 변환",    9,   False, C_BLUE, "ko-KR"),
        ("    ├── checker.py  ← 페이지 상태 검증",                  9,   False, C_BLUE, "ko-KR"),
        ("    ├── registry.py ← ItemDef · CRF 항목 메타데이터",     9,   False, C_BLUE, "ko-KR"),
        ("    └── reporter.py ← Excel / HTML 결과 리포트 생성",     9,   False, C_BLUE, "ko-KR"),
    ]
    paras = [(t, sz, b, col, "l", lang) for (t, sz, b, col, lang) in tree_lines]
    add_rect(slide, Inches(0.42), Inches(1.48), Inches(6.1), Inches(5.67),
             fill=C_CARD, border=C_BORDER, radius=8000)
    add_textbox(slide, Inches(0.55), Inches(1.55), Inches(5.9), Inches(5.55),
                paras, vert_anchor="top",
                margin_l=Inches(0), margin_t=Inches(0.05),
                margin_r=Inches(0), margin_b=Inches(0))

    # ── 우: 역할 카드 ──
    rx = Inches(6.88)
    rw = Inches(6.0)

    draw_section_header(slide, rx, Inches(1.05), rw, "핵심 클래스 역할 분담")

    draw_card(slide, rx, Inches(1.48), rw, Inches(2.25), 1,
              "CDMAgent — 브라우저 조작 레이어",
              [
                  ("Chrome extension ↔ Daemon ↔ Python 통신을 담당합니다.", 9, False, C_DARK, "ko-KR"),
                  ("", 8, False, C_DARK, "ko-KR"),
                  ("inspect()  현재 페이지 스냅샷 반환 (PageSnapshot)", 8.5, False, C_GRAY, "ko-KR"),
                  ("set_date / set_text / select_radio / select_option", 8.5, False, C_GRAY, "en-US"),
                  ("click_save / click_save_next / go_to_page / go_back", 8.5, False, C_GRAY, "en-US"),
                  ("has_query / check_result  DVS Query 발생 여부 검증", 8.5, False, C_GRAY, "ko-KR"),
                  ("list_pages()  CRF 페이지 코드 목록 표로 출력", 8.5, False, C_GRAY, "ko-KR"),
              ])

    draw_card(slide, rx, Inches(3.82), rw, Inches(3.28), 2,
              "DVSRunner — DVS 검증 오케스트레이터",
              [
                  ("Excel DVS 파일을 읽어 자동 검증 후 결과를 저장합니다.", 9, False, C_DARK, "ko-KR"),
                  ("", 8, False, C_DARK, "ko-KR"),
                  ("① DVSRunner(agent, excel_path)  초기화", 8.5, False, C_GRAY, "ko-KR"),
                  ("② runner.dry_run()  Excel 파싱 결과 미리보기 (브라우저 X)", 8.5, False, C_GRAY, "ko-KR"),
                  ("③ runner.run_all()  전체 자동 실행 + Excel/HTML 리포트", 8.5, False, C_GRAY, "ko-KR"),
                  ("   또는 runner.run_one(dvs_id)  단건 실행", 8.5, False, C_GRAY, "ko-KR"),
                  ("④ runner.record() + flush_to_excel()  수동 결과 기록", 8.5, False, C_GRAY, "ko-KR"),
                  ("", 8, False, C_DARK, "ko-KR"),
                  ("결과: R열(Result) · S열(Comment) · 색상 코딩 자동 적용", 8.5, False, C_GREEN, "ko-KR"),
                  ("PASS=초록  FAIL=빨강  SKIP·ERROR=노랑  원본 덮어쓰기 없음", 8.5, False, C_GRAY, "ko-KR"),
              ])


# ── 슬라이드 3: 효과 및 향후 ──────────────────────────────────────────────────

def build_slide3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    draw_title_bar(slide,
                   "자동화 도입 효과 및 향후 개발 방안",
                   "수동 DM 업무 자동화 전환 · 품질 향상 · 확장 로드맵")

    lx = Inches(0.42)
    lw = Inches(5.9)
    rx = Inches(6.85)
    rw = Inches(6.05)

    draw_section_header(slide, lx, Inches(1.05), lw, "자동화 도입 효과")
    draw_section_header(slide, rx, Inches(1.05), rw, "향후 개발 방안")

    # ── 효과 카드 3개 ──
    draw_card(slide, lx, Inches(1.48), lw, Inches(1.70), 1,
              "반복 작업 자동화 → 업무 효율 대폭 향상",
              [
                  ("CRF 입력 · DVS 검증 시 케이스마다 수행하던 클릭 · 복붙 · 결과", 9, False, C_DARK, "ko-KR"),
                  ("기록 작업을 완전히 자동화합니다.", 9, False, C_DARK, "ko-KR"),
                  ("", 8, False, C_DARK, "ko-KR"),
                  ("케이스당 5~10분 수동 작업 → 수십 초 자동 처리", 8.5, True, C_GREEN, "ko-KR"),
              ])

    draw_card(slide, lx, Inches(3.26), lw, Inches(1.70), 2,
              "인적 오류 감소 → 데이터 품질 향상",
              [
                  ("입력 오류 · 결과 오기록 등 인적 실수를 방지하고", 9, False, C_DARK, "ko-KR"),
                  ("check_result()로 Query 발생 여부를 즉시 자동 검증합니다.", 9, False, C_DARK, "ko-KR"),
                  ("", 8, False, C_DARK, "ko-KR"),
                  ("PASS/FAIL 색상 코딩 자동 적용 → 결과 검토 직관성 향상", 8.5, True, C_BLUE, "ko-KR"),
              ])

    draw_card(slide, lx, Inches(5.04), lw, Inches(1.70), 3,
              "Jupyter 기반 → 범용성 및 재현성 확보",
              [
                  ("Jupyter Notebook에서 HTML 카드로 결과를 즉시 확인하고", 9, False, C_DARK, "ko-KR"),
                  ("스크립트를 저장해 동일 작업을 언제든 재현 가능합니다.", 9, False, C_DARK, "ko-KR"),
                  ("", 8, False, C_DARK, "ko-KR"),
                  ("Python 생태계(pandas, openpyxl 등)와 자연스럽게 연동", 8.5, False, C_GRAY, "ko-KR"),
              ])

    # ── 향후 개발 카드 3개 ──
    draw_card(slide, rx, Inches(1.48), rw, Inches(1.70), 1,
              "멀티-스터디 / 배치 실행 지원",
              [
                  ("여러 스터디를 순차 또는 병렬로 자동 실행하는 Batch 모드를", 9, False, C_DARK, "ko-KR"),
                  ("추가해 대규모 DVS 처리 효율을 높입니다.", 9, False, C_DARK, "ko-KR"),
                  ("", 8, False, C_DARK, "ko-KR"),
                  ("runner.run_all(studies=[...]) 형태의 API 확장 예정", 8.5, False, C_GRAY, "en-US"),
              ])

    draw_card(slide, rx, Inches(3.26), rw, Inches(1.70), 2,
              "AI 기반 CRF 이상 탐지 연동",
              [
                  ("PageSnapshot 데이터를 LLM에 전달해 이상 값을 자동 감지하고", 9, False, C_DARK, "ko-KR"),
                  ("Query 예측 모델을 사전 적용합니다.", 9, False, C_DARK, "ko-KR"),
                  ("", 8, False, C_DARK, "ko-KR"),
                  ("Claude API 연동 → DVS 스펙 자동 생성 파이프라인 구축", 8.5, False, C_GRAY, "ko-KR"),
              ])

    draw_card(slide, rx, Inches(5.04), rw, Inches(1.70), 3,
              "결과 대시보드 및 보고서 자동화",
              [
                  ("DVS 검증 결과를 웹 대시보드로 시각화하고", 9, False, C_DARK, "ko-KR"),
                  ("스터디별 통계를 자동 집계해 PM/Sponsor와 공유합니다.", 9, False, C_DARK, "ko-KR"),
                  ("", 8, False, C_DARK, "ko-KR"),
                  ("HTML/Excel 리포트 → 공유 워크플로우 완전 자동화", 8.5, True, C_ORANGE, "ko-KR"),
              ])


# ── 메인 ──────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    build_slide1(prs)
    build_slide2(prs)
    build_slide3(prs)

    if os.path.exists(OUTPUT):
        os.remove(OUTPUT)
    prs.save(OUTPUT)
    size_mb = os.path.getsize(OUTPUT) / 1024 / 1024
    print(f"저장 완료: {OUTPUT}")
    print(f"파일 크기: {size_mb:.1f} MB  |  슬라이드: {len(prs.slides)}장")


if __name__ == "__main__":
    main()
